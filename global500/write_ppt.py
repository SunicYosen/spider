
import os
import time
from pptx           import Presentation
from pptx.util      import Cm, Pt
from pptx.enum.text import PP_ALIGN

from get_data_global500 import get_data_global500
from get_industry_cat   import get_industry_cat
from get_canvas         import get_canvas

def write_ppt(pic_path = 'output', industry_dict={}, company_list=[], template='template.pptx'):  
    prs = Presentation(template)
    time_begin = time.time()
    index = 1

    for key in industry_dict:
        '''
         添加分类页
        '''
        # 添加文本框
        slide   = prs.slides.add_slide(prs.slide_layouts[1])
        textbox = slide.shapes.add_textbox( left=Cm(2),
                                            top=Cm(8),
                                            width=Cm(20),
                                            height=Cm(4))
        # 向文本框加入文字
        tf             = textbox.text_frame
        para           = tf.add_paragraph()    # 添加段落
        para.text      = industry_dict[key]
        para.alignment = PP_ALIGN.CENTER       # 居中
        # 设置字体
        font           = para.font
        font.size      = Pt(64)       # 大小
        # font.name    = '华文彩云'    # 字体
        font.bold      = True         # 加粗
        # font.italic  = True         # 倾斜
        # font.color.rgb = RGBColor(225, 225, 0)  # 黄色

        for company in company_list:
            duration  = round(time.time() - time_begin, 2)
            remaining = round(duration * 500 / (index) - duration, 2)

            print("[+] Getting Data :{}/{}，Cost Time: {}s， Remainimg: {}s".format(index, 500, duration, remaining), end="\r")

            # 按类别分
            if company.industry == industry_dict[key]:
                index     = index + 1
                slide = prs.slides.add_slide(prs.slide_layouts[1])

                title_shape = slide.shapes.title
                title_shape.text = "{}\n{}".format(company.company_name, company.company_name_en)

                img_name1 = '{}/{}/{}_{}.png'.format(pic_path, company.rank, company.company_name, company.image1_suffix).replace('&','')
                img_name2 = '{}/{}/{}_{}.png'.format(pic_path, company.rank, company.company_name, company.image2_suffix).replace('&','')
                img_name3 = '{}/{}/{}_{}.png'.format(pic_path, company.rank, company.company_name, company.image3_suffix).replace('&','')

                if not os.path.exists(os.path.join(str(pic_path), str(company.rank))):
                    os.makedirs(os.path.join(str(pic_path), str(company.rank)))

                if not (os.path.exists(img_name1) & os.path.exists(img_name2) & os.path.exists(img_name3)):
                    get_canvas(company.data_site, os.path.join(str(pic_path), str(company.rank)))

                try:
                    slide.shapes.add_picture(image_file=img_name1,
                                            left=Cm(0.5),
                                            top=Cm(5),
                                            width=Cm(11),
                                            height=Cm(10))
                except:
                    print("[-]: {}:{} does not found!".format(company.company_name, img_name1))

                try:
                    slide.shapes.add_picture(image_file=img_name2,
                                            left=Cm(11.5),
                                            top=Cm(5),
                                            width=Cm(11),
                                            height=Cm(10))
                except:
                    print("[-]: {}:{} does not found!".format(company.company_name, img_name2))
                    
                try:
                    slide.shapes.add_picture(image_file=img_name3,
                                            left=Cm(22.5),
                                            top=Cm(5),
                                            width=Cm(11),
                                            height=Cm(10))
                except:
                    print("[-]: {}:{} does not found!".format(company.company_name, img_name3))

        prs.save('result.pptx')

    print("[+] Info: Write PPTX Done!")

if __name__ == '__main__':
    root_url       = 'https://www.caifuzhongwen.com/fortune500/paiming/global500/2021_%e4%b8%96%e7%95%8c500%e5%bc%ba.htm'
    industry_dict  = get_industry_cat(root_url)
    company_list   = get_data_global500(root_url)
    write_ppt(industry_dict=industry_dict, company_list=company_list)
